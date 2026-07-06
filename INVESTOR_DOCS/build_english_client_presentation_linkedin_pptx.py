from __future__ import annotations

import math
import tempfile
from pathlib import Path
from typing import Iterable

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
OUT_DIR = BASE_DIR / "presentation_output"
OUT_DIR.mkdir(exist_ok=True)
OUT_FILE = OUT_DIR / "CADRO_HR_Modul_Anlatimli_Presentation_LinkedIn_v3.pptx"
LOGO_PATH = BASE_DIR / "assets" / "cadro_logo_extracted.png"

SLIDE_FONT = "Aptos SemiBold"

SHOT_ROOT = Path("/Users/turgaybozkus/Desktop/cadro-app-screenshots-webp/english")
MAIN_DIR = SHOT_ROOT / "main"
STATE_DIR = SHOT_ROOT / "states"
LINKEDIN_DIR = BASE_DIR / "assets" / "linkedin_posts"

TMP_DIR = Path(tempfile.mkdtemp(prefix="cadro_pptx_img_"))


NAVY = RGBColor(21, 33, 58)
INK = RGBColor(28, 34, 48)
MUTED = RGBColor(91, 100, 120)
BRAND = RGBColor(13, 158, 223)
GREEN = RGBColor(15, 157, 88)


MODULE_STORY: list[dict[str, str]] = [
    {
        "title": "Employee Management",
        "image": "main/employees.webp",
        "what": "Maintains a clean, centralized employee master record.",
        "fix": "Fixes duplicate personnel data and manual record drift.",
        "impact": "Improves payroll readiness and HR data reliability.",
    },
    {
        "title": "Onboarding & Offboarding",
        "image": "main/onboarding.webp",
        "what": "Runs onboarding lifecycle with structured checklists and readiness status.",
        "fix": "Fixes first-week process gaps and inconsistent handovers.",
        "impact": "Accelerates time-to-productivity for new hires.",
    },
    {
        "title": "Applicant Tracking (ATS)",
        "image": "states/ats--new-candidate-modal.webp",
        "what": "Visualizes candidates across hiring stages in one pipeline.",
        "fix": "Fixes candidate loss and non-transparent hiring progression.",
        "impact": "Shortens time-to-hire and improves hiring quality.",
    },
    {
        "title": "Assets & Inventory",
        "image": "states/assets--new-asset-modal.webp",
        "what": "Tracks assignment, stock, and return state for company assets.",
        "fix": "Fixes ownership ambiguity and missing handover records.",
        "impact": "Reduces asset loss and replacement cost.",
    },
    {
        "title": "Purchase Requests",
        "image": "states/purchase-requests--new-purchase-request-modal.webp",
        "what": "Standardizes request-to-approval purchasing workflow.",
        "fix": "Fixes untracked procurement and uncontrolled spend requests.",
        "impact": "Improves budget governance and approval speed.",
    },
    {
        "title": "Corporate Request Forms",
        "image": "states/request-forms--new-request-modal.webp",
        "what": "Digitizes internal request forms with clear workflow ownership.",
        "fix": "Fixes free-format request chaos and missed follow-ups.",
        "impact": "Increases response consistency and process traceability.",
    },
    {
        "title": "Training (LMS)",
        "image": "states/training--new-training-modal.webp",
        "what": "Organizes trainings, participants, schedules, and completion actions.",
        "fix": "Fixes training visibility gaps across teams.",
        "impact": "Improves capability development and compliance readiness.",
    },
    {
        "title": "Knowledge Base & Policies",
        "image": "states/knowledge-base--new-content-modal.webp",
        "what": "Publishes policies and tracks acknowledgement status.",
        "fix": "Fixes policy communication inconsistency and low policy visibility.",
        "impact": "Strengthens governance and policy adoption.",
    },
    {
        "title": "E-Dossier",
        "image": "states/e-dossier--compliance-tab.webp",
        "what": "Manages digital employee document archive and approval actions.",
        "fix": "Fixes scattered document storage and weak compliance control.",
        "impact": "Reduces audit risk and speeds compliance checks.",
    },
    {
        "title": "Locations & Sites",
        "image": "states/locations--add-site-modal.webp",
        "what": "Controls location setup for site-based workforce operations.",
        "fix": "Fixes location ambiguity in distributed teams.",
        "impact": "Improves site-level coordination and attendance integrity.",
    },
    {
        "title": "Support Messages",
        "image": "states/support-center--new-ticket-modal.webp",
        "what": "Captures support issues via structured ticket submission.",
        "fix": "Fixes informal support requests lost in chats and emails.",
        "impact": "Improves response quality and accountability.",
    },
]


HR_CHALLENGE_SLIDES: list[dict[str, object]] = [
    {
        "title": "HR Challenge 1: Fragmented Employee Data",
        "subtitle": "When HR records are spread across files, emails, and disconnected tools",
        "bullets": [
            "Duplicate employee records create payroll inconsistencies and reporting confusion",
            "Onboarding and offboarding steps are missed because ownership is not centralized",
            "Management cannot trust one single source of truth for workforce decisions",
        ],
    },
    {
        "title": "HR Challenge 2: Slow and Opaque Workflows",
        "subtitle": "When approvals and requests depend on chats and manual follow-ups",
        "bullets": [
            "Recruitment, purchase, and internal requests stall due to unclear approval paths",
            "Support and operational requests are lost between channels",
            "Cycle times increase while accountability decreases across teams",
        ],
    },
    {
        "title": "HR Challenge 3: Compliance and Capability Gaps",
        "subtitle": "When policy, document, and training governance is inconsistent",
        "bullets": [
            "Policy acknowledgement and document control are difficult to audit",
            "Training visibility is limited, making compliance readiness weaker",
            "Distributed locations operate with inconsistent standards and low transparency",
        ],
    },
]


def to_png(src: Path) -> Path:
    dst = TMP_DIR / f"{src.stem}.png"
    if dst.exists():
        return dst
    with PILImage.open(src) as img:
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGB")
        elif img.mode == "RGBA":
            bg = PILImage.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            img = bg
        img.save(dst, format="PNG")
    return dst


def add_title(slide, text: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Calibri"

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.45))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED
        sp.font.name = "Calibri"


def add_bullet_slide(prs: Presentation, title: str, subtitle: str, bullets: Iterable[str], footer: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(4.7))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = INK
        p.font.name = "Calibri"
        p.space_after = Pt(10)

    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35))
        ff = footer_box.text_frame
        ff.clear()
        fp = ff.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(10)
        fp.font.color.rgb = MUTED
        fp.font.name = "Calibri"
        fp.alignment = PP_ALIGN.RIGHT


def add_image_grid_slide(prs: Presentation, title: str, subtitle: str, images: list[Path], cols: int, rows: int, top: float = 1.45):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)

    left = 0.5
    width = 12.33
    gap_x = 0.15
    gap_y = 0.15

    cell_w = (width - gap_x * (cols - 1)) / cols
    cell_h = (7.1 - top - gap_y * (rows - 1)) / rows

    for idx, image_path in enumerate(images[: cols * rows]):
        r = idx // cols
        c = idx % cols
        x = left + c * (cell_w + gap_x)
        y = top + r * (cell_h + gap_y)
        png = to_png(image_path)
        slide.shapes.add_picture(str(png), Inches(x), Inches(y), width=Inches(cell_w), height=Inches(cell_h))


def add_module_story_slide(prs: Presentation, module: dict[str, str], image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Başlık: tam genişlik, yatayda ortalı, Aptos SemiBold, tema rengi ──
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(10.1), Inches(0.75))
    tf = title_box.text_frame
    tf.clear()
    tp = tf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    run = tp.add_run()
    run.text = module["title"]          # "Step:" öneki yok
    run.font.name = SLIDE_FONT
    run.font.bold = True
    # Renk ayarlanmıyor → PowerPoint tema rengi (siyah/koyu)

    # ── CADRO logosu sağ üst ──
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), Inches(10.88), Inches(0.06),
            width=Inches(2.45), height=Inches(1.64)
        )

    # ── Ekran görüntüsü sol taraf ──
    png = to_png(image_path)
    slide.shapes.add_picture(str(png), Inches(0.55), Inches(1.6), width=Inches(7.1), height=Inches(4.55))

    # ── Açıklama metin kutusu sağ taraf ──
    txt = slide.shapes.add_textbox(Inches(7.83), Inches(2.2), Inches(4.95), Inches(3.8))
    tf = txt.text_frame
    tf.word_wrap = True
    tf.clear()

    lines = [
        ("What it does", module["what"]),
        ("What it fixes", module["fix"]),
        ("Business impact", module["impact"]),
    ]

    for idx, (label, value) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.JUSTIFY
        run = p.add_run()
        run.text = f"{label}: {value}"
        run.font.name = SLIDE_FONT
        run.font.bold = True
        run.font.size = Pt(20)
        # Renk ayarlanmıyor → PowerPoint tema rengi
        if idx < len(lines) - 1:
            p.space_after = Pt(12)


def pick_module_images() -> list[tuple[dict[str, str], Path]]:
    all_images = sorted(
        [p for p in MAIN_DIR.glob("*.webp") if p.name != ".DS_Store"]
        + [p for p in STATE_DIR.glob("*.webp") if p.name != ".DS_Store"]
    )
    index = {f"{p.parent.name}/{p.name}".lower(): p for p in all_images}
    selected: list[tuple[dict[str, str], Path]] = []
    for module in MODULE_STORY:
        key = module["image"].lower()
        if key in index:
            selected.append((module, index[key]))
    return selected


def chunks(items: list[Path], size: int):
    for i in range(0, len(items), size):
        yield items[i : i + size]


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    main_images = sorted([p for p in MAIN_DIR.glob("*.webp") if p.name != ".DS_Store"])
    state_images = sorted([p for p in STATE_DIR.glob("*.webp") if p.name != ".DS_Store"])
    linkedin_images = sorted([p for p in LINKEDIN_DIR.glob("*.jpg") if p.name != ".DS_Store"])

    add_bullet_slide(
        prs,
        "CADRO HR",
        "Step-by-Step HR Transformation Story",
        [
            "Focused module walk-through with only high-impact screens",
            "Each step explains: what the module does, what HR problem it fixes, and business impact",
            "Public proof from LinkedIn posts + product interface proof",
            "IP-safe narrative for enterprise clients and investors",
        ],
        "CADRO HR | Client-facing | English",
    )

    for challenge in HR_CHALLENGE_SLIDES:
        add_bullet_slide(
            prs,
            str(challenge["title"]),
            str(challenge["subtitle"]),
            challenge["bullets"],
        )

    add_bullet_slide(
        prs,
        "How CADRO HRMS Solves These Challenges",
        "Problem-to-module resolution flow",
        [
            "Challenge 1 -> Employee Management, Onboarding & Offboarding, E-Dossier",
            "Challenge 2 -> ATS, Purchase Requests, Corporate Request Forms, Support Messages",
            "Challenge 3 -> Training (LMS), Knowledge Base & Policies, Locations & Sites",
            "Result -> Standardized execution, traceable operations, faster HR decisions",
        ],
    )

    module_pairs = pick_module_images()
    for module, image_path in module_pairs:
        add_module_story_slide(prs, module, image_path)

    if linkedin_images:
        per_slide = 4
        total = math.ceil(len(linkedin_images) / per_slide)
        for idx, group in enumerate(chunks(linkedin_images, per_slide), 1):
            add_image_grid_slide(
                prs,
                f"LinkedIn Product Posts ({idx}/{total})",
                "External proof: public product communication and market visibility",
                group,
                cols=2,
                rows=2,
                top=1.55,
            )

    key_states = [
        p
        for p in state_images
        if p.name
        in {
            "performance--okr-tab.webp",
            "employees--new-record-modal.webp",
            "attendance--payroll-report-tab.webp",
            "purchase-requests--new-purchase-request-modal.webp",
        }
    ]
    if key_states:
        add_image_grid_slide(
            prs,
            "Key Workflow Proof",
            "Critical state transitions that show real operational depth",
            key_states,
            cols=2,
            rows=2,
            top=1.55,
        )

    add_bullet_slide(
        prs,
        "ROI and Commercial Case",
        "Why this is financially and strategically meaningful",
        [
            "Base scenario payback is around 10.2 months with measurable process gains",
            "Aggressive scenario payback can approach 5.2 months",
            "Subscription-led model with module upsell increases lifetime value",
            "Multilingual readiness (TR / EN / AR) supports regional expansion",
            "Public digital presence (cadro.io + LinkedIn) strengthens market credibility",
        ],
    )

    cta = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(cta, "Next Step", "Pilot rollout, commercial due diligence, and scale plan")
    box = cta.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.8), Inches(2.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CADRO is ready to improve HR execution quickly while keeping core implementation IP protected."
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    prs.save(str(OUT_FILE))
    print(f"Created {OUT_FILE}")
    print(f"LinkedIn images: {len(linkedin_images)}")
    print(f"Main images: {len(main_images)}")
    print(f"State images: {len(state_images)}")


if __name__ == "__main__":
    build()
